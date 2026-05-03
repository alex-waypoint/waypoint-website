#!/usr/bin/env python3
"""
Waypoint Client Presentation Template Generator
Run: python3 generate_deck.py
Output: waypoint_client_template.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ═══════════════════════════════════════════════════
# BRAND TOKENS
# ═══════════════════════════════════════════════════
DARK        = RGBColor(0x09, 0x0B, 0x0F)
DARK_2      = RGBColor(0x11, 0x15, 0x20)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEAL        = RGBColor(0x0F, 0x5E, 0x7A)
TEAL_H      = RGBColor(0x1A, 0x70, 0x90)
TEAL_LIGHT  = RGBColor(0x5A, 0xC4, 0xDE)
TEAL_PALE   = RGBColor(0xDF, 0xF0, 0xF5)
T1          = RGBColor(0x09, 0x0B, 0x0F)
T2          = RGBColor(0x40, 0x40, 0x48)
T3          = RGBColor(0x8A, 0x8A, 0x92)
BG          = RGBColor(0xF8, 0xF8, 0xF6)
BG_2        = RGBColor(0xEF, 0xEF, 0xED)
BORDER      = RGBColor(0xD4, 0xD4, 0xD2)
GREEN       = RGBColor(0x4A, 0xDE, 0x80)

# ─── Layout constants ───────────────────────────────
W   = Inches(10)
H   = Inches(5.625)
M   = Inches(0.55)     # margin
CW  = Inches(8.9)      # content width
MID = W / 2

# ═══════════════════════════════════════════════════
# PRIMITIVE HELPERS
# ═══════════════════════════════════════════════════

def mk_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def add_slide(prs):
    blank = prs.slide_layouts[6]
    return prs.slides.add_slide(blank)

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, fc=None, lc=None, lw=Pt(1)):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.shadow.inherit = False
    if fc is not None:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fc
    else:
        sh.fill.background()
    if lc is not None:
        sh.line.color.rgb = lc
        sh.line.width = lw
    else:
        sh.line.fill.background()
    return sh

def tb(slide, text, l, t, w, h,
       sz=Pt(12), co=T1, bold=False, italic=False,
       align=PP_ALIGN.LEFT, font="Inter", wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text       = text
    run.font.name  = font
    run.font.size  = sz
    run.font.color.rgb = co
    run.font.bold  = bold
    run.font.italic = italic
    return txb, tf

def para(tf, text, sz=Pt(11), co=T1, bold=False,
         align=PP_ALIGN.LEFT, sb=Pt(4), font="Inter"):
    p   = tf.add_paragraph()
    p.alignment    = align
    p.space_before = sb
    run = p.add_run()
    run.text       = text
    run.font.name  = font
    run.font.size  = sz
    run.font.color.rgb = co
    run.font.bold  = bold
    return p

# ─── Reusable components ────────────────────────────

def add_footer(slide, dark=True):
    co = WHITE if dark else T3
    tb(slide, "WAYPOINT",
       M, H - Inches(0.38), Inches(1.5), Inches(0.3),
       sz=Pt(7.5), co=co, bold=True)
    tb(slide, "Confidential · Not for Distribution",
       W - M - Inches(2.8), H - Inches(0.38), Inches(2.8), Inches(0.3),
       sz=Pt(7.5), co=co, align=PP_ALIGN.RIGHT)

def add_eyebrow(slide, text, l, t, dark=False):
    co = TEAL_LIGHT if dark else TEAL
    tb(slide, text, l, t, Inches(5), Inches(0.28),
       sz=Pt(7.5), co=co, bold=True)

def add_divider_line(slide, y, dark=True):
    co = TEAL if dark else BORDER
    box(slide, M, y, CW, Inches(0.025), fc=co)

def add_corner_marks(slide, dark=True):
    """Geospatial corner bracket marks"""
    co = RGBColor(0xFF,0xFF,0xFF) if dark else BORDER
    lw = Pt(1.2)
    # top-left
    box(slide, M, Inches(0.3), Inches(0.22), Inches(0.015), fc=co)
    box(slide, M, Inches(0.3), Inches(0.015), Inches(0.22), fc=co)
    # bottom-right
    box(slide, W - M - Inches(0.22), H - Inches(0.52), Inches(0.22), Inches(0.015), fc=co)
    box(slide, W - M - Inches(0.015), H - Inches(0.52), Inches(0.015), Inches(0.22), fc=co)

def add_section_tag(slide, num, label, dark=True):
    """Small section number tag for section dividers"""
    co_num  = TEAL_LIGHT if dark else TEAL
    co_lbl  = RGBColor(0xFF,0xFF,0xFF) if dark else T3
    tb(slide, f"0{num}", M, H/2 - Inches(1.1), Inches(1), Inches(0.8),
       sz=Pt(72), co=co_num, bold=True, font="Inter")
    tb(slide, label.upper(), M, H/2 + Inches(0.0), Inches(6), Inches(0.35),
       sz=Pt(8), co=co_lbl, bold=True, font="Inter")

# ═══════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════

# ─── 01 · COVER ─────────────────────────────────────
def slide_cover(prs):
    s = add_slide(prs)
    set_bg(s, DARK)

    # Right accent panel
    box(s, W - Inches(2.6), Inches(0), Inches(2.6), H, fc=DARK_2)
    box(s, W - Inches(2.6), Inches(0), Inches(0.025), H, fc=TEAL)

    # Top teal bar
    box(s, Inches(0), Inches(0), W, Inches(0.03), fc=TEAL)

    # Wordmark
    tb(s, "WAYPOINT", M, Inches(0.42), Inches(2), Inches(0.4),
       sz=Pt(14), co=WHITE, bold=True)

    # Eyebrow
    add_eyebrow(s, "GEOSPATIAL DATA GENERATION & INTELLIGENCE", M, Inches(1.45), dark=True)

    # Main headline placeholder
    tb(s, "Prepared for [CLIENT NAME]",
       M, Inches(1.85), Inches(6.5), Inches(1.1),
       sz=Pt(34), co=WHITE, bold=True)

    # Subtitle placeholder
    tb(s, "[Engagement Title or Project Description]",
       M, Inches(3.05), Inches(6.5), Inches(0.45),
       sz=Pt(14), co=T3)

    # Divider line
    box(s, M, Inches(3.6), Inches(3.5), Inches(0.02), fc=TEAL)

    # Meta line
    _,tf = tb(s, "[Month YYYY]",
       M, Inches(3.75), Inches(4), Inches(0.3),
       sz=Pt(10), co=T3)
    para(tf, "waypoint-geo.com  ·  a.declaro@creattura.com",
         sz=Pt(10), co=T3, sb=Pt(2))

    # Right panel subtle text
    tb(s, "CONFIDENTIAL", W - Inches(2.4), H - Inches(0.5), Inches(2), Inches(0.25),
       sz=Pt(7), co=T3, align=PP_ALIGN.CENTER)

    add_corner_marks(s, dark=True)
    add_footer(s, dark=True)
    return s

# ─── 02 · TABLE OF CONTENTS ─────────────────────────
def slide_toc(prs):
    s = add_slide(prs)
    set_bg(s, BG)

    # Left accent bar
    box(s, M, Inches(0.5), Inches(0.04), Inches(4.5), fc=TEAL)

    # Title
    tb(s, "Contents", M + Inches(0.22), Inches(0.5), Inches(3), Inches(0.55),
       sz=Pt(28), co=T1, bold=True)

    sections = [
        ("01", "About Waypoint"),
        ("02", "Capabilities"),
        ("03", "Why Human Annotation"),
        ("04", "Our Process"),
        ("05", "Industry Applications"),
        ("06", "Engagement & Pricing"),
        ("07", "Case Study"),
        ("08", "Next Steps"),
    ]

    col_w = Inches(4.0)
    for i, (num, label) in enumerate(sections):
        col = i // 4
        row = i % 4
        lx = M + col * Inches(4.6)
        ty = Inches(1.35) + row * Inches(0.82)

        # Row background
        box(s, lx, ty, col_w, Inches(0.62),
            fc=BG_2, lc=BORDER, lw=Pt(0.75))

        # Number
        tb(s, num, lx + Inches(0.18), ty + Inches(0.12),
           Inches(0.45), Inches(0.38),
           sz=Pt(11), co=TEAL, bold=True)

        # Label
        tb(s, label, lx + Inches(0.65), ty + Inches(0.12),
           Inches(3.2), Inches(0.38),
           sz=Pt(12), co=T1, bold=False)

    add_footer(s, dark=False)
    return s

# ─── 03 · SECTION DIVIDER ───────────────────────────
def slide_divider(prs, num, title, subtitle=""):
    s = add_slide(prs)
    set_bg(s, DARK_2)

    # Left teal accent band
    box(s, Inches(0), Inches(0), Inches(0.08), H, fc=TEAL)

    # Large background number (decorative)
    tb(s, f"0{num}", Inches(0.5), H/2 - Inches(1.5), Inches(3), Inches(2.2),
       sz=Pt(160), co=RGBColor(0x14,0x1A,0x28), bold=True, align=PP_ALIGN.LEFT)

    # Section label
    add_eyebrow(s, f"SECTION 0{num}", Inches(0.5), H/2 - Inches(0.38), dark=True)

    # Title
    tb(s, title, Inches(0.5), H/2 - Inches(0.05), Inches(7), Inches(0.7),
       sz=Pt(36), co=WHITE, bold=True)

    # Subtitle
    if subtitle:
        tb(s, subtitle, Inches(0.5), H/2 + Inches(0.72), Inches(6.5), Inches(0.45),
           sz=Pt(13), co=T3)

    # Right corner decoration
    box(s, W - Inches(1.8), Inches(1.2), Inches(1.8), Inches(0.025), fc=TEAL)
    box(s, W - Inches(0.025), Inches(1.2), Inches(0.025), Inches(2.5),
        fc=RGBColor(0x14,0x1A,0x28))

    add_corner_marks(s, dark=True)
    add_footer(s, dark=True)
    return s

# ─── 04 · COMPANY OVERVIEW ──────────────────────────
def slide_overview(prs):
    s = add_slide(prs)
    set_bg(s, BG)

    add_eyebrow(s, "ABOUT WAYPOINT", M, Inches(0.45))
    tb(s, "Geospatial Data Generation\n& Intelligence",
       M, Inches(0.72), Inches(5.5), Inches(0.95),
       sz=Pt(24), co=T1, bold=True)
    add_divider_line(s, Inches(1.72), dark=False)

    # Left: description
    _, tf = tb(s, "Waypoint is a geospatial data generation and intelligence company powering the organizations that map, monitor, and manage the physical world.",
               M, Inches(1.85), Inches(5.1), Inches(0.75),
               sz=Pt(12), co=T2, wrap=True)

    _, tf2 = tb(s, "We combine trained human teams, structured QA, and documented workflows to deliver annotation, classification, and analysis outputs that are precise, audit-ready, and GeoAI-compatible.",
                M, Inches(2.68), Inches(5.1), Inches(0.8),
                sz=Pt(12), co=T2, wrap=True)

    # Key facts - right column
    facts = [
        ("Human-validated", "Every dataset delivered through documented human QA"),
        ("GeoAI-ready",     "Structured outputs compatible with ML training pipelines"),
        ("Multi-domain",    "Agriculture, infrastructure, land cover, defense"),
    ]
    for i, (label, desc) in enumerate(facts):
        ty = Inches(1.78) + i * Inches(1.05)
        box(s, Inches(5.6), ty, Inches(3.85), Inches(0.9),
            fc=WHITE, lc=BORDER, lw=Pt(0.75))
        box(s, Inches(5.6), ty, Inches(0.04), Inches(0.9), fc=TEAL)
        tb(s, label, Inches(5.76), ty + Inches(0.1), Inches(3.5), Inches(0.3),
           sz=Pt(11), co=T1, bold=True)
        tb(s, desc, Inches(5.76), ty + Inches(0.42), Inches(3.5), Inches(0.38),
           sz=Pt(9.5), co=T3, wrap=True)

    add_footer(s, dark=False)
    return s

# ─── 05 · THE CHALLENGE ─────────────────────────────
def slide_challenge(prs):
    s = add_slide(prs)
    set_bg(s, BG)

    add_eyebrow(s, "THE CHALLENGE", M, Inches(0.45))
    tb(s, "The gap between raw geospatial data\nand actionable intelligence is where projects fail.",
       M, Inches(0.72), Inches(8), Inches(0.95),
       sz=Pt(22), co=T1, bold=True, wrap=True)
    add_divider_line(s, Inches(1.74), dark=False)

    pain_points = [
        ("Inconsistent Annotation Quality",
         "Manual workflows without documented QA produce datasets that fail model validation — costly rework, delayed timelines.",
         "Quality"),
        ("Fragmented Geospatial Pipelines",
         "Organizations collect data from multiple sources (satellite, LiDAR, field) but lack the integration layer to produce coherent intelligence outputs.",
         "Integration"),
        ("Automated Tools Can't Solve Context",
         "GeoAI models require high-quality training data and human-in-the-loop validation that automated tools alone cannot provide at the required precision.",
         "Precision"),
    ]
    card_w = Inches(2.82)
    for i, (title, body, tag) in enumerate(pain_points):
        lx = M + i * Inches(3.02)
        ty = Inches(1.92)
        box(s, lx, ty, card_w, Inches(2.9),
            fc=WHITE, lc=BORDER, lw=Pt(0.75))
        # Top teal bar
        box(s, lx, ty, card_w, Inches(0.04), fc=TEAL)
        # Tag
        box(s, lx + Inches(0.18), ty + Inches(0.18),
            Inches(1.0), Inches(0.24), fc=TEAL_PALE)
        tb(s, tag.upper(), lx + Inches(0.22), ty + Inches(0.2),
           Inches(0.9), Inches(0.22),
           sz=Pt(7), co=TEAL, bold=True)
        # Title
        tb(s, title, lx + Inches(0.18), ty + Inches(0.56),
           card_w - Inches(0.36), Inches(0.55),
           sz=Pt(12), co=T1, bold=True, wrap=True)
        # Body
        tb(s, body, lx + Inches(0.18), ty + Inches(1.22),
           card_w - Inches(0.36), Inches(1.48),
           sz=Pt(10), co=T2, wrap=True)

    add_footer(s, dark=False)
    return s

# ─── 06 · DATA GENERATION CAPABILITY ────────────────
def slide_capability_data(prs):
    s = add_slide(prs)
    set_bg(s, BG)

    # Left content (55%)
    add_eyebrow(s, "CAPABILITY 01", M, Inches(0.45))
    tb(s, "Data Generation", M, Inches(0.72), Inches(5), Inches(0.55),
       sz=Pt(28), co=T1, bold=True)
    add_divider_line(s, Inches(1.35), dark=False)

    _, tf = tb(s, "From raw imagery to structured, annotation-complete datasets — processed through trained human teams and documented QA workflows.",
               M, Inches(1.5), Inches(4.9), Inches(0.72),
               sz=Pt(11.5), co=T2, wrap=True)

    bullets = [
        "Feature extraction — buildings, roads, boundaries, vegetation",
        "Object detection & classification at scale",
        "Semantic and instance segmentation",
        "Change detection across temporal datasets",
        "Multi-source fusion: satellite, aerial, LiDAR, ground truth",
        "GeoAI training dataset production",
    ]
    _, tf_b = tb(s, bullets[0],
                 M, Inches(2.38), Inches(4.9), Inches(0.28),
                 sz=Pt(10.5), co=T2)
    for b in bullets[1:]:
        para(tf_b, b, sz=Pt(10.5), co=T2, sb=Pt(5))

    # Right: visual placeholder panel
    box(s, Inches(5.6), Inches(0.5), Inches(3.85), Inches(4.68),
        fc=DARK, lc=TEAL, lw=Pt(0.75))
    tb(s, "[ DATA GENERATION\nVISUALIZATION ]",
       Inches(5.6), Inches(2.2), Inches(3.85), Inches(1.2),
       sz=Pt(12), co=T3, align=PP_ALIGN.CENTER, wrap=True)
    # Corner marks on placeholder
    box(s, Inches(5.68), Inches(0.58), Inches(0.3), Inches(0.015), fc=TEAL)
    box(s, Inches(5.68), Inches(0.58), Inches(0.015), Inches(0.3), fc=TEAL)
    box(s, Inches(9.37), Inches(5.1), Inches(0.08), Inches(0.015), fc=TEAL)
    box(s, Inches(9.44), Inches(5.02), Inches(0.015), Inches(0.08), fc=TEAL)

    add_footer(s, dark=False)
    return s

# ─── 07 · GEO INTELLIGENCE CAPABILITY ───────────────
def slide_capability_geo(prs):
    s = add_slide(prs)
    set_bg(s, BG)

    add_eyebrow(s, "CAPABILITY 02", M, Inches(0.45))
    tb(s, "Geospatial Intelligence", M, Inches(0.72), Inches(5.5), Inches(0.55),
       sz=Pt(28), co=T1, bold=True)
    add_divider_line(s, Inches(1.35), dark=False)

    _, tf = tb(s, "From processed datasets to verified intelligence outputs — through spatial analysis, model validation, and human-in-the-loop synthesis.",
               M, Inches(1.5), Inches(4.9), Inches(0.72),
               sz=Pt(11.5), co=T2, wrap=True)

    bullets = [
        "GeoAI model validation & accuracy assessment",
        "Human-in-the-loop inference review",
        "Spatial analysis & pattern detection",
        "Environmental & infrastructure monitoring",
        "Change detection analysis (multi-temporal)",
        "Situational awareness & decision support datasets",
    ]
    _, tf_b = tb(s, bullets[0],
                 M, Inches(2.38), Inches(4.9), Inches(0.28),
                 sz=Pt(10.5), co=T2)
    for b in bullets[1:]:
        para(tf_b, b, sz=Pt(10.5), co=T2, sb=Pt(5))

    # Right: monitoring layer slideshow placeholder
    box(s, Inches(5.6), Inches(0.5), Inches(3.85), Inches(4.68),
        fc=DARK, lc=TEAL, lw=Pt(0.75))
    tb(s, "[ GEOSPATIAL INTELLIGENCE\nVISUALIZATION ]",
       Inches(5.6), Inches(2.2), Inches(3.85), Inches(1.2),
       sz=Pt(12), co=T3, align=PP_ALIGN.CENTER, wrap=True)
    box(s, Inches(5.68), Inches(0.58), Inches(0.3), Inches(0.015), fc=TEAL)
    box(s, Inches(5.68), Inches(0.58), Inches(0.015), Inches(0.3), fc=TEAL)

    add_footer(s, dark=False)
    return s

# ─── 08 · WHY HUMAN ANNOTATION ──────────────────────
def slide_why_human(prs):
    s = add_slide(prs)
    set_bg(s, WHITE)

    # Teal header band
    box(s, Inches(0), Inches(0), W, Inches(1.65), fc=DARK)

    add_eyebrow(s, "DIFFERENTIATION", M, Inches(0.35), dark=True)
    tb(s, "Why Human Annotation Matters",
       M, Inches(0.65), Inches(8), Inches(0.7),
       sz=Pt(26), co=WHITE, bold=True)

    reasons = [
        ("Contextual judgment",
         "Geospatial features require interpretation: a structure is not just a polygon — it is a greenhouse, a reservoir, a military asset. Context requires human cognition."),
        ("Edge case resolution",
         "Automated models fail at boundaries, occlusion, and ambiguity. Human annotators resolve these with documented decision rules, not probabilistic guesses."),
        ("Accountable QA chain",
         "Human validation creates an audit trail. Every label has a reviewer. Every dataset has a QA record. Automated pipelines cannot replicate this accountability structure."),
    ]
    card_w = Inches(2.82)
    for i, (title, body) in enumerate(reasons):
        lx = M + i * Inches(3.02)
        ty = Inches(1.82)
        box(s, lx, ty, card_w, Inches(3.22), fc=BG, lc=BORDER, lw=Pt(0.75))
        # Number
        tb(s, f"0{i+1}", lx + Inches(0.18), ty + Inches(0.15),
           Inches(0.5), Inches(0.45),
           sz=Pt(24), co=TEAL, bold=True)
        # Title
        tb(s, title, lx + Inches(0.18), ty + Inches(0.68),
           card_w - Inches(0.36), Inches(0.4),
           sz=Pt(12), co=T1, bold=True)
        # Body
        tb(s, body, lx + Inches(0.18), ty + Inches(1.15),
           card_w - Inches(0.36), Inches(1.9),
           sz=Pt(10.5), co=T2, wrap=True)

    add_footer(s, dark=False)
    return s

# ─── 09 · QA / QUALITY FRAMEWORK ────────────────────
def slide_qa(prs):
    s = add_slide(prs)
    set_bg(s, BG)

    add_eyebrow(s, "QUALITY FRAMEWORK", M, Inches(0.45))
    tb(s, "Structured QA at Every Stage",
       M, Inches(0.72), Inches(6), Inches(0.55),
       sz=Pt(26), co=T1, bold=True)
    add_divider_line(s, Inches(1.35), dark=False)

    # QA Stage table
    stages = [
        ("Intake & Specification",
         "Project brief, data specs, annotation guidelines documented before any work begins.",
         "Prevents downstream ambiguity"),
        ("Annotator Execution",
         "Trained teams apply documented rules. Every annotator is domain-specialized.",
         "Domain expertise per project"),
        ("Peer Review Layer",
         "Independent review of 100% of annotations against specification.",
         "Two-person accountability"),
        ("QA Lead Audit",
         "Senior QA lead audits random sample + edge cases. Issues logged and resolved.",
         "Statistical + judgment review"),
        ("Client Delivery Package",
         "Dataset delivered with QA report, accuracy metrics, and decision log.",
         "Full audit trail"),
    ]
    for i, (stage, desc, tag) in enumerate(stages):
        ty = Inches(1.52) + i * Inches(0.72)
        # Row background, alternating
        row_bg = WHITE if i % 2 == 0 else BG_2
        box(s, M, ty, CW, Inches(0.68), fc=row_bg, lc=BORDER, lw=Pt(0.5))
        # Step number
        box(s, M, ty, Inches(0.4), Inches(0.68), fc=TEAL)
        tb(s, str(i+1), M + Inches(0.1), ty + Inches(0.17),
           Inches(0.2), Inches(0.34),
           sz=Pt(13), co=WHITE, bold=True, align=PP_ALIGN.CENTER)
        # Stage name
        tb(s, stage, M + Inches(0.55), ty + Inches(0.12),
           Inches(2.4), Inches(0.44),
           sz=Pt(11), co=T1, bold=True)
        # Description
        tb(s, desc, M + Inches(3.1), ty + Inches(0.12),
           Inches(4.1), Inches(0.44),
           sz=Pt(10), co=T2, wrap=True)
        # Tag
        box(s, M + Inches(7.35), ty + Inches(0.17),
            Inches(1.4), Inches(0.26), fc=TEAL_PALE)
        tb(s, tag, M + Inches(7.4), ty + Inches(0.19),
           Inches(1.3), Inches(0.24),
           sz=Pt(8), co=TEAL, bold=True, wrap=True)

    add_footer(s, dark=False)
    return s

# ─── 10 · HOW WE WORK (PROCESS) ─────────────────────
def slide_process(prs):
    s = add_slide(prs)
    set_bg(s, BG)

    add_eyebrow(s, "OUR PROCESS", M, Inches(0.45))
    tb(s, "From Brief to Delivery",
       M, Inches(0.72), Inches(6), Inches(0.55),
       sz=Pt(26), co=T1, bold=True)
    add_divider_line(s, Inches(1.35), dark=False)

    steps = [
        ("Discovery",       "Project brief, data specs, scope alignment"),
        ("Team Setup",      "Domain-matched annotators, guidelines finalized"),
        ("Execution",       "Annotation + real-time QA monitoring"),
        ("Review & QA",     "Peer review, QA audit, edge case resolution"),
        ("Delivery",        "Dataset + QA report + decision log"),
    ]
    step_w = Inches(1.68)
    step_h = Inches(2.6)
    for i, (title, desc) in enumerate(steps):
        lx = M + i * (step_w + Inches(0.1))
        ty = Inches(1.7)

        # Card
        box(s, lx, ty, step_w, step_h, fc=WHITE, lc=BORDER, lw=Pt(0.75))

        # Number circle (simulated with small teal box)
        box(s, lx + Inches(0.16), ty + Inches(0.18),
            Inches(0.38), Inches(0.38), fc=TEAL)
        tb(s, str(i+1), lx + Inches(0.22), ty + Inches(0.2),
           Inches(0.26), Inches(0.3),
           sz=Pt(12), co=WHITE, bold=True, align=PP_ALIGN.CENTER)

        # Title
        tb(s, title, lx + Inches(0.12), ty + Inches(0.72),
           step_w - Inches(0.24), Inches(0.38),
           sz=Pt(12), co=T1, bold=True)

        # Desc
        tb(s, desc, lx + Inches(0.12), ty + Inches(1.18),
           step_w - Inches(0.24), Inches(1.25),
           sz=Pt(9.5), co=T2, wrap=True)

        # Arrow between steps
        if i < len(steps) - 1:
            ax = lx + step_w + Inches(0.02)
            tb(s, "→", ax, ty + step_h/2 - Inches(0.22),
               Inches(0.12), Inches(0.3),
               sz=Pt(14), co=TEAL, align=PP_ALIGN.CENTER)

    add_footer(s, dark=False)
    return s

# ─── 11 · INDUSTRY APPLICATIONS ─────────────────────
def slide_industries(prs):
    s = add_slide(prs)
    set_bg(s, BG)

    add_eyebrow(s, "INDUSTRY APPLICATIONS", M, Inches(0.45))
    tb(s, "Where Waypoint Operates",
       M, Inches(0.72), Inches(6), Inches(0.55),
       sz=Pt(26), co=T1, bold=True)
    add_divider_line(s, Inches(1.35), dark=False)

    industries = [
        ("Agriculture & Land Use",
         "Field boundary mapping, crop classification, vegetation monitoring, yield prediction datasets.",
         "01"),
        ("Infrastructure & Utilities",
         "Road network extraction, building footprints, pipeline mapping, change detection for maintenance.",
         "02"),
        ("Environmental Monitoring",
         "Land cover classification, deforestation detection, wetlands mapping, carbon stock estimation.",
         "03"),
        ("Defense & Intelligence",
         "Object detection, facility classification, terrain analysis, activity monitoring.",
         "04"),
    ]
    card_w = Inches(4.3)
    for i, (title, desc, num) in enumerate(industries):
        col = i % 2
        row = i // 2
        lx  = M + col * Inches(4.5)
        ty  = Inches(1.58) + row * Inches(1.6)
        box(s, lx, ty, card_w, Inches(1.42), fc=WHITE, lc=BORDER, lw=Pt(0.75))
        box(s, lx, ty, Inches(0.04), Inches(1.42), fc=TEAL)
        # Num tag
        tb(s, num, lx + Inches(0.18), ty + Inches(0.15),
           Inches(0.4), Inches(0.3),
           sz=Pt(10), co=TEAL, bold=True)
        tb(s, title, lx + Inches(0.62), ty + Inches(0.15),
           card_w - Inches(0.8), Inches(0.3),
           sz=Pt(12), co=T1, bold=True)
        tb(s, desc, lx + Inches(0.18), ty + Inches(0.55),
           card_w - Inches(0.36), Inches(0.76),
           sz=Pt(10), co=T2, wrap=True)

    add_footer(s, dark=False)
    return s

# ─── 12 · ENGAGEMENT MODEL ───────────────────────────
def slide_engagement(prs):
    s = add_slide(prs)
    set_bg(s, BG)

    add_eyebrow(s, "ENGAGEMENT MODEL", M, Inches(0.45))
    tb(s, "How We Work Together",
       M, Inches(0.72), Inches(6), Inches(0.55),
       sz=Pt(26), co=T1, bold=True)
    add_divider_line(s, Inches(1.35), dark=False)

    models = [
        ("Project-Based",
         "Defined scope, timeline, and deliverables. Ideal for one-time dataset production or pilot engagements.",
         "Fixed Scope  ·  Fixed Price",
         TEAL_PALE),
        ("Ongoing Retainer",
         "Recurring annotation capacity with dedicated teams. Ideal for organizations with continuous geospatial data pipelines.",
         "Monthly Volume  ·  Dedicated Team",
         BG_2),
        ("Embedded Partnership",
         "Waypoint team integrated into client workflows. Ideal for enterprise GeoAI programs requiring consistent human-in-the-loop coverage.",
         "Custom SLA  ·  Full Integration",
         TEAL_PALE),
    ]
    card_w = Inches(2.82)
    for i, (title, desc, tag, bg_c) in enumerate(models):
        lx = M + i * Inches(3.02)
        ty = Inches(1.62)
        box(s, lx, ty, card_w, Inches(3.3), fc=bg_c, lc=BORDER, lw=Pt(0.75))
        box(s, lx, ty, card_w, Inches(0.05), fc=TEAL)
        tb(s, title, lx + Inches(0.18), ty + Inches(0.2),
           card_w - Inches(0.36), Inches(0.38),
           sz=Pt(13), co=T1, bold=True)
        tb(s, desc, lx + Inches(0.18), ty + Inches(0.7),
           card_w - Inches(0.36), Inches(1.5),
           sz=Pt(10.5), co=T2, wrap=True)
        # Tag line
        box(s, lx + Inches(0.18), ty + Inches(2.58),
            card_w - Inches(0.36), Inches(0.28), fc=WHITE)
        tb(s, tag, lx + Inches(0.22), ty + Inches(2.6),
           card_w - Inches(0.44), Inches(0.26),
           sz=Pt(8), co=TEAL, bold=True)

    # CTA note
    tb(s, "All engagements begin with a scoping consultation — no commitment required.",
       M, Inches(5.1), CW, Inches(0.28),
       sz=Pt(10), co=T3, align=PP_ALIGN.CENTER)

    add_footer(s, dark=False)
    return s

# ─── 13 · CASE STUDY ────────────────────────────────
def slide_case_study(prs):
    s = add_slide(prs)
    set_bg(s, BG)

    # Dark header strip
    box(s, Inches(0), Inches(0), W, Inches(1.35), fc=DARK)

    add_eyebrow(s, "CASE STUDY", M, Inches(0.3), dark=True)
    tb(s, "[Project Title — Client Industry]",
       M, Inches(0.6), Inches(8), Inches(0.5),
       sz=Pt(22), co=WHITE, bold=True)

    # Three-panel structure
    panels = [
        ("Challenge",
         "[Describe the client's core geospatial data problem — what they needed, why existing approaches were insufficient, and the scale or complexity involved.]\n\n[2–3 sentences]"),
        ("Approach",
         "[Describe how Waypoint scoped and executed the engagement — teams deployed, methodology, QA stages, timeline, and any specialized techniques used.]\n\n[2–3 sentences]"),
        ("Outcome",
         "[Describe the delivered dataset, the measurable quality metrics, and the downstream impact for the client — model performance, operational decisions enabled, etc.]\n\n[2–3 sentences]"),
    ]
    panel_w = Inches(2.82)
    for i, (title, body) in enumerate(panels):
        lx = M + i * Inches(3.02)
        ty = Inches(1.48)
        box(s, lx, ty, panel_w, Inches(3.68), fc=WHITE, lc=BORDER, lw=Pt(0.75))
        # Header
        box(s, lx, ty, panel_w, Inches(0.4), fc=BG_2)
        tb(s, title.upper(), lx + Inches(0.18), ty + Inches(0.08),
           panel_w - Inches(0.36), Inches(0.28),
           sz=Pt(8), co=TEAL, bold=True)
        tb(s, body, lx + Inches(0.18), ty + Inches(0.55),
           panel_w - Inches(0.36), Inches(3.0),
           sz=Pt(10), co=T2, wrap=True)

    # Metrics bar
    box(s, M, Inches(5.18), CW, Inches(0.025), fc=BORDER)
    metrics = ["[Accuracy %]", "[Annotations Delivered]", "[Timeline]", "[Client Sector]"]
    mw = CW / len(metrics)
    for i, m in enumerate(metrics):
        lx = M + i * mw
        tb(s, m, lx + Inches(0.1), Inches(5.22), mw - Inches(0.2), Inches(0.28),
           sz=Pt(9), co=T3, bold=False)

    add_footer(s, dark=False)
    return s

# ─── 14 · NEXT STEPS ────────────────────────────────
def slide_next_steps(prs):
    s = add_slide(prs)
    set_bg(s, DARK)

    # Top accent
    box(s, Inches(0), Inches(0), W, Inches(0.03), fc=TEAL)

    add_eyebrow(s, "MOVING FORWARD", M, Inches(0.52), dark=True)
    tb(s, "Three steps to get started",
       M, Inches(0.82), Inches(6), Inches(0.65),
       sz=Pt(30), co=WHITE, bold=True)

    steps = [
        ("01", "Scoping Call",
         "30-minute call to align on project requirements, data types, annotation specifications, and timeline expectations."),
        ("02", "Pilot Engagement",
         "A small-batch pilot dataset to validate methodology, team fit, and quality standards — before full-scale commitment."),
        ("03", "Full Engagement",
         "Structured delivery against agreed specifications, with QA reporting at every stage."),
    ]
    step_w = Inches(2.82)
    for i, (num, title, desc) in enumerate(steps):
        lx = M + i * Inches(3.02)
        ty = Inches(1.7)
        box(s, lx, ty, step_w, Inches(2.98),
            fc=DARK_2, lc=RGBColor(0x1E,0x24,0x30), lw=Pt(0.75))
        box(s, lx, ty, step_w, Inches(0.04), fc=TEAL)

        # Large number
        tb(s, num, lx + Inches(0.18), ty + Inches(0.14),
           Inches(0.6), Inches(0.55),
           sz=Pt(30), co=TEAL, bold=True)
        tb(s, title, lx + Inches(0.18), ty + Inches(0.78),
           step_w - Inches(0.36), Inches(0.38),
           sz=Pt(13), co=WHITE, bold=True)
        tb(s, desc, lx + Inches(0.18), ty + Inches(1.28),
           step_w - Inches(0.36), Inches(1.5),
           sz=Pt(10.5), co=T3, wrap=True)

    add_corner_marks(s, dark=True)
    add_footer(s, dark=True)
    return s

# ─── 15 · CONTACT & CLOSE ───────────────────────────
def slide_close(prs):
    s = add_slide(prs)
    set_bg(s, DARK_2)

    # Left teal accent
    box(s, Inches(0), Inches(0), Inches(0.06), H, fc=TEAL)

    # Top teal bar
    box(s, Inches(0), Inches(0), W, Inches(0.03), fc=TEAL)

    # Wordmark large
    tb(s, "WAYPOINT", M, Inches(1.4), Inches(5), Inches(0.8),
       sz=Pt(52), co=WHITE, bold=True)

    tb(s, "Geospatial Data Generation & Intelligence",
       M, Inches(2.28), Inches(6), Inches(0.4),
       sz=Pt(14), co=TEAL_LIGHT)

    # Divider
    box(s, M, Inches(2.82), Inches(4.5), Inches(0.02), fc=TEAL)

    # Contact
    _, tf = tb(s, "a.declaro@creattura.com",
               M, Inches(3.05), Inches(4), Inches(0.35),
               sz=Pt(13), co=T3)
    para(tf, "waypoint-geo.com", sz=Pt(13), co=T3, sb=Pt(3))
    para(tf, "[Phone / Calendar link]", sz=Pt(13), co=T3, sb=Pt(3))

    # Right side CTA
    box(s, Inches(6.2), Inches(2.0), Inches(3.25), Inches(1.45), fc=TEAL)
    tb(s, "Book a\nConsultation",
       Inches(6.45), Inches(2.12), Inches(2.7), Inches(1.1),
       sz=Pt(22), co=WHITE, bold=True, wrap=True)
    tb(s, "No commitment required",
       Inches(6.45), Inches(3.52), Inches(2.7), Inches(0.3),
       sz=Pt(9.5), co=TEAL_PALE)

    add_corner_marks(s, dark=True)
    add_footer(s, dark=True)
    return s

# ─── 16 · APPENDIX DIVIDER ──────────────────────────
def slide_appendix(prs):
    s = add_slide(prs)
    set_bg(s, BG)
    box(s, Inches(0), Inches(0), W, Inches(0.04), fc=BORDER)

    add_eyebrow(s, "APPENDIX", M, Inches(0.52))
    tb(s, "Supporting Materials",
       M, Inches(0.82), Inches(6), Inches(0.55),
       sz=Pt(28), co=T1, bold=True)
    add_divider_line(s, Inches(1.42), dark=False)

    appendix_items = [
        "Team Credentials & Domain Specializations",
        "Sample QA Report Format",
        "Data Specification Template",
        "Annotation Guideline Sample",
        "Pricing Framework",
        "References & Data Sources",
    ]
    _, tf = tb(s, appendix_items[0],
               M, Inches(1.65), CW, Inches(0.32),
               sz=Pt(12), co=T2)
    for item in appendix_items[1:]:
        para(tf, item, sz=Pt(12), co=T2, sb=Pt(8))

    add_footer(s, dark=False)
    return s


# ═══════════════════════════════════════════════════
# ASSEMBLE DECK
# ═══════════════════════════════════════════════════

def build():
    prs = mk_prs()

    slide_cover(prs)                                           # 01 Cover
    slide_toc(prs)                                             # 02 Table of Contents

    slide_divider(prs, 1, "About Waypoint",                    # 03
                  "Who we are and what we do")
    slide_overview(prs)                                        # 04 Company Overview
    slide_challenge(prs)                                       # 05 The Challenge

    slide_divider(prs, 2, "Capabilities",                      # 06
                  "What Waypoint delivers")
    slide_capability_data(prs)                                 # 07 Data Generation
    slide_capability_geo(prs)                                  # 08 Geo Intelligence
    slide_why_human(prs)                                       # 09 Why Human
    slide_qa(prs)                                              # 10 QA Framework

    slide_divider(prs, 3, "Our Process",                       # 11
                  "How we scope, execute and deliver")
    slide_process(prs)                                         # 12 Process Flow

    slide_divider(prs, 4, "Industry Applications",             # 13
                  "Sectors we serve")
    slide_industries(prs)                                      # 14 Industries

    slide_divider(prs, 5, "Engagement",                        # 15
                  "How to work with Waypoint")
    slide_engagement(prs)                                      # 16 Engagement Model
    slide_case_study(prs)                                      # 17 Case Study

    slide_next_steps(prs)                                      # 18 Next Steps
    slide_close(prs)                                           # 19 Contact & Close
    slide_appendix(prs)                                        # 20 Appendix

    out = "/Users/alexisdeclaro/Desktop/waypoint/claude-website/waypoint_client_template.pptx"
    prs.save(out)
    print(f"✓ Saved: {out}")
    print(f"  Slides: {len(prs.slides)}")

if __name__ == "__main__":
    build()
